"""
Office Toolkit - 办公文件工具箱
================================
功能：
1. 文件分类整理（按类型/日期/关键词）
2. 办公模板库（日报、周报、会议纪要、合同等）
3. 文件格式转换（TXT/DOCX/图片等）
"""

import os, re, json, shutil, hashlib, time, datetime, subprocess, zipfile
from pathlib import Path
from flask import Flask, request, jsonify, send_file, send_from_directory, render_template_string

app = Flask(__name__, static_folder="static", static_url_path="/static")

# ==================== CONFIG ====================
TOOLKIT_DIR = Path(__file__).parent
TEMPLATES_DIR = TOOLKIT_DIR / "templates"
CONVERTERS_DIR = TOOLKIT_DIR / "converters"
OUTPUT_DIR = TOOLKIT_DIR / "output"
OUTPUT_DIR.mkdir(exist_ok=True)

# ==================== FILE FILER ====================
FILE_CATEGORIES = {
    "文档": [".doc", ".docx", ".txt", ".pdf", ".rtf", ".md", ".wps"],
    "表格": [".xls", ".xlsx", ".csv", ".et"],
    "演示": [".ppt", ".pptx", ".pps", ".dps"],
    "图片": [".jpg", ".jpeg", ".png", ".gif", ".bmp", ".tiff", ".webp", ".svg"],
    "视频": [".mp4", ".avi", ".mkv", ".mov", ".wmv", ".flv"],
    "音频": [".mp3", ".wav", ".flac", ".aac", ".ogg", ".wma"],
    "压缩包": [".zip", ".rar", ".7z", ".tar", ".gz"],
    "代码": [".py", ".js", ".html", ".css", ".java", ".cpp", ".c", ".h", 
             ".ts", ".go", ".rs", ".php", ".sql", ".json", ".xml", ".yaml", ".yml"],
    "设计": [".psd", ".ai", ".cdr", ".sketch", ".fig"],
    "邮件": [".eml", ".msg", ".pst"],
    "其他": [],
}

def auto_classify(file_path):
    """按扩展名自动分类文件"""
    ext = Path(file_path).suffix.lower()
    for category, exts in FILE_CATEGORIES.items():
        if ext in exts:
            return category
    return "其他"

@app.route("/api/classify", methods=["POST"])
def api_classify():
    """
    分类指定目录中的文件
    Body: {"dir": "C:/path/to/files", "mode": "type|date|keyword", "keyword": "xxx"}
    """
    data = request.get_json() or {}
    target_dir = data.get("dir", "")
    mode = data.get("mode", "type")  # type/date/keyword
    keyword = data.get("keyword", "")

    if not target_dir or not os.path.isdir(target_dir):
        return jsonify({"ok": False, "message": "无效目录"}), 400

    files = []
    for f in os.listdir(target_dir):
        fp = os.path.join(target_dir, f)
        if os.path.isfile(fp):
            stat = os.stat(fp)
            mtime = datetime.datetime.fromtimestamp(stat.st_mtime)
            files.append({
                "name": f, "path": fp, "size": stat.st_size,
                "mtime": mtime.strftime("%Y-%m-%d %H:%M"),
                "ext": os.path.splitext(f)[1].lower(),
            })

    if not files:
        return jsonify({"ok": True, "files": [], "categories": {}})

    if mode == "type":
        # 按文件类型归类
        grouped = {}
        for f in files:
            cat = auto_classify(f["name"])
            if cat not in grouped:
                grouped[cat] = []
            grouped[cat].append(f)

        # 预览归类结果（不实际移动）
        result = {}
        for cat, items in grouped.items():
            result[cat] = {
                "count": len(items),
                "total_size": sum(i["size"] for i in items),
                "files": [i["name"] for i in items],
            }

        return jsonify({"ok": True, "files": files, "categories": result, "mode": mode})

    elif mode == "date":
        # 按月归类
        grouped = {}
        for f in files:
            month_key = f["mtime"][:7]  # YYYY-MM
            if month_key not in grouped:
                grouped[month_key] = []
            grouped[month_key].append(f)
        result = {}
        for mk, items in sorted(grouped.items()):
            result[mk] = {
                "count": len(items),
                "total_size": sum(i["size"] for i in items),
                "files": [i["name"] for i in items],
            }
        return jsonify({"ok": True, "files": files, "categories": result, "mode": mode})

    elif mode == "keyword" and keyword:
        # 按关键词筛选
        matched = [f for f in files if keyword.lower() in f["name"].lower()]
        result = {
            "匹配结果": {
                "count": len(matched),
                "total_size": sum(i["size"] for i in matched),
                "files": [i["name"] for i in matched],
            }
        }
        return jsonify({"ok": True, "files": files, "categories": result, "mode": mode})

    return jsonify({"ok": False, "message": "参数错误"}), 400


@app.route("/api/classify_execute", methods=["POST"])
def api_classify_execute():
    """实际执行文件移动归类"""
    data = request.get_json() or {}
    target_dir = data.get("dir", "")
    mode = data.get("mode", "type")

    if not target_dir or not os.path.isdir(target_dir):
        return jsonify({"ok": False, "message": "无效目录"}), 400

    moved = 0
    errors = []
    target_path = Path(target_dir)

    for f in os.listdir(target_dir):
        fp = target_path / f
        if not fp.is_file():
            continue

        if mode == "type":
            cat = auto_classify(f)
        elif mode == "date":
            mtime = datetime.datetime.fromtimestamp(fp.stat().st_mtime)
            cat = mtime.strftime("%Y-%m")
        else:
            continue

        cat_dir = target_path / cat
        cat_dir.mkdir(exist_ok=True)

        # 避免同名冲突
        dest = cat_dir / f
        if dest.exists():
            base, ext = os.path.splitext(f)
            dest = cat_dir / f"{base}_{int(time.time())}{ext}"

        try:
            shutil.move(str(fp), str(dest))
            moved += 1
        except Exception as e:
            errors.append({"file": f, "error": str(e)})

    return jsonify({"ok": True, "moved": moved, "errors": errors})


# ==================== OFFICE TEMPLATES ====================
# 内置模板定义
BUILTIN_TEMPLATES = [
    {
        "id": "daily_report",
        "title": "日报模板",
        "category": "汇报",
        "desc": "标准日报格式，含今日工作、明日计划、问题反馈",
        "format": "txt",
        "content": """\
=== 工作日报 ===
日期：{date}
姓名：{name}

一、今日完成工作
1. 
2. 
3. 

二、明日工作计划
1. 
2. 
3. 

三、遇到的问题与建议
1. 
2. 

四、需要协调事项
""",
    },
    {
        "id": "weekly_report",
        "title": "周报模板",
        "category": "汇报",
        "desc": "标准周报格式，含本周总结、下周计划、关键指标",
        "format": "txt",
        "content": """\
=== 工作周报 ===
周期：{week_range}
姓名：{name}

一、本周重点工作
| 序号 | 工作内容 | 完成情况 | 备注 |
|------|---------|---------|------|
| 1    |         |         |      |
| 2    |         |         |      |
| 3    |         |         |      |

二、关键指标
1. 
2. 

三、下周工作计划
1. 
2. 
3. 

四、需要支持的事项
""",
    },
    {
        "id": "meeting_minutes",
        "title": "会议纪要模板",
        "category": "会议",
        "desc": "标准会议纪要，含议题、决议、待办事项",
        "format": "txt",
        "content": """\
=== 会议纪要 ===
会议主题：{topic}
日期时间：{date}
地点：{location}
主持人：{host}
参会人员：{attendees}
记录人：{recorder}

一、会议议程
1. 
2. 
3. 

二、主要讨论内容
1. 
2. 
3. 

三、决议事项
1. 
2. 

四、待办事项
| 事项 | 负责人 | 截止日期 |
|------|--------|---------|
|      |        |         |
|      |        |         |

五、下次会议安排
""",
    },
    {
        "id": "contract_brief",
        "title": "合同模板（简易版）",
        "category": "合同",
        "desc": "简易合作协议/合同范本",
        "format": "txt",
        "content": """\
=== 合作协议 ===
协议编号：{contract_no}
签订日期：{date}

甲方：{party_a}
乙方：{party_b}

一、合作内容
1. 

二、合作期限
自 {start_date} 至 {end_date}

三、费用及支付方式
1. 
2. 

四、双方权利与义务
1. 
2. 
3. 

五、违约责任
1. 

六、争议解决
本协议履行过程中发生争议，双方应友好协商解决；协商不成的，提交甲方所在地人民法院诉讼解决。

七、其他
本协议一式两份，甲乙双方各执一份，具有同等法律效力。

甲方（盖章）：________    乙方（盖章）：________
签字：________              签字：________
日期：________              日期：________
""",
    },
    {
        "id": "leave_request",
        "title": "请假/调休申请",
        "category": "人事",
        "desc": "标准请假审批单，含事假/病假/年假/调休",
        "format": "txt",
        "content": """\
=== 请假申请单 ===
申请人：{name}
部门：{dept}
日期：{date}

请假类型：□事假 □病假 □年假 □调休 □婚假 □其他

请假时间：自 {start_date} {start_time} 至 {end_date} {end_time}
共计：{days} 天

请假事由：
1. 

工作交接人：{handover}

部门负责人意见：________
签字：________  日期：________
""",
    },
    {
        "id": "project_plan",
        "title": "项目计划模板",
        "category": "项目",
        "desc": "项目里程碑计划和进度表",
        "format": "txt",
        "content": """\
=== 项目计划书 ===
项目名称：{project_name}
项目经理：{manager}
起止日期：{start} — {end}

一、项目目标
1. 

二、里程碑节点
| 阶段 | 开始日 | 结束日 | 交付物 | 负责人 |
|------|--------|--------|--------|--------|
|      |        |        |        |        |
|      |        |        |        |        |

三、资源需求
1. 人员：
2. 预算：
3. 设备：

四、风险评估
1. 
2. 

五、沟通计划
1. 
""",
    },
    {
        "id": "okr_template",
        "title": "OKR目标模板",
        "category": "管理",
        "desc": "OKR（目标和关键结果）制定模板",
        "format": "txt",
        "content": """\
=== OKR 目标设定 ===
周期：{period}
姓名/团队：{name}

Objective 1（目标1）：
{objective_1}

Key Results：
KR1: {kr_1_1}
KR2: {kr_1_2}
KR3: {kr_1_3}

---
Objective 2（目标2）：
{objective_2}

Key Results：
KR1: {kr_2_1}
KR2: {kr_2_2}
KR3: {kr_2_3}

---
信心指数评分：
- 5分：极有把握
- 3分：中等把握
- 1分：挑战很大
""",
    },
    {
        "id": "personal_summary",
        "title": "个人工作总结",
        "category": "汇报",
        "desc": "季度/年度个人工作总结报告",
        "format": "txt",
        "content": """\
=== {year}年{period}个人工作总结 ===
姓名：{name}
岗位：{position}

一、工作概况
1. 

二、重点工作成果
1. {achievement_1}
2. {achievement_2}
3. {achievement_3}

三、能力提升与学习
1. 
2. 

四、存在的不足
1. 
2. 

五、下阶段工作计划
1. 
2. 
3. 

六、个人感悟
""",
    },
    {
        "id": "todo_list",
        "title": "待办事项清单",
        "category": "日常",
        "desc": "简洁的待办事项列表模板",
        "format": "txt",
        "content": """\
=== 待办事项清单 ===
日期：{date}

优先级：🔴高 🟡中 🟢低

🔴 重要且紧急
□ 1. 
□ 2. 

🟡 重要不紧急
□ 1. 
□ 2. 

🟢 日常事务
□ 1. 
□ 2. 

备注：
""",
    },
    {
        "id": "meeting_agenda",
        "title": "会议议程模板",
        "category": "会议",
        "desc": "标准的会议议程安排表",
        "format": "txt",
        "content": """\
=== 会议议程 ===
会议名称：{meeting_name}
日期：{date}
时间：{start_time} - {end_time}

议程安排：
---
| 时间 | 议题 | 主讲人 | 时长 |
|------|------|--------|------|
|      | 开场   |        | 5min |
|      |       |        |      |
|      |       |        |      |
|      | 自由讨论 |      |      |
|      | 总结   |        | 5min |

备注：
""",
    },
    {
        "id": "travel_request",
        "title": "出差申请单",
        "category": "人事",
        "desc": "公司出差审批单模板",
        "format": "txt",
        "content": """\
=== 出差申请单 ===
申请人：{name}
部门：{dept}
日期：{date}

出差地点：{destination}
出差事由：{reason}
同行人员：{companions}

起止时间：{start_date} — {end_date}

预计费用：
- 交通：____元
- 住宿：____元
- 餐补：____元
- 其他：____元
- 合计：____元

部门审批：________
财务审批：________
总经理审批：________
""",
    },
    {
        "id": "training_report",
        "title": "培训记录模板",
        "category": "人事",
        "desc": "员工培训签到和培训记录表",
        "format": "txt",
        "content": """\
=== 培训记录 ===
培训主题：{topic}
培训讲师：{trainer}
培训日期：{date}
培训地点：{location}
培训时长：{duration} 小时

参训人员名单：
1. 
2. 
3. 

培训内容：
1. 
2. 
3. 

培训效果评估（满意度 1-5分）：
- 内容相关性：____分
- 讲师表现：____分
- 实用性：____分
- 平均分：____分

备注：
""",
    },
    {
        "id": "work_log",
        "title": "工作日志模板",
        "category": "日常",
        "desc": "简洁的工作时间记录日志",
        "format": "txt",
        "content": """\
=== 工作日志 ===
日期：{date}

| 时间 | 工作内容 | 耗时 | 备注 |
|------|---------|------|------|
| 09:00-10:00 |       | 1h   |      |
| 10:00-11:00 |       | 1h   |      |
| 11:00-12:00 |       | 1h   |      |
| 12:00-13:30 | 午休   |      |      |
| 13:30-14:30 |       | 1h   |      |
| 14:30-15:30 |       | 1h   |      |
| 15:30-16:30 |       | 1h   |      |
| 16:30-17:30 |       | 1h   |      |
| 17:30-18:00 |       | 0.5h |      |

总计有效工时：____ 小时
今日反思：
""",
    },
    {
        "id": "invoice_record",
        "title": "报销单据模板",
        "category": "财务",
        "desc": "费用报销清单明细表",
        "format": "txt",
        "content": """\
=== 费用报销单 ===
申请人：{name}
部门：{dept}
申请日期：{date}
报销事由：{reason}

费用明细：
| 序号 | 费用类别 | 金额(元) | 票据编号 | 备注 |
|------|---------|---------|---------|------|
| 1    |         |         |         |      |
| 2    |         |         |         |      |
| 3    |         |         |         |      |
| 合计 |         | 0.00    |         |      |

大写金额：________

审批流程：
部门负责人：________  日期：________
财务审核：________    日期：________
""",
    },
]

@app.route("/api/templates")
def api_templates():
    """获取所有模板列表"""
    category = request.args.get("category", "")
    templates = BUILTIN_TEMPLATES
    if category:
        templates = [t for t in templates if t["category"] == category]
    return jsonify({
        "ok": True,
        "templates": [{"id": t["id"], "title": t["title"],
                       "category": t["category"], "desc": t["desc"],
                       "format": t["format"]} for t in templates],
        "categories": list(set(t["category"] for t in BUILTIN_TEMPLATES)),
    })

@app.route("/api/templates/<template_id>")
def api_template_detail(template_id):
    """获取模板详细内容"""
    for t in BUILTIN_TEMPLATES:
        if t["id"] == template_id:
            return jsonify({"ok": True, **t})
    return jsonify({"ok": False, "message": "模板不存在"}), 404

@app.route("/api/templates/generate", methods=["POST"])
def api_template_generate():
    """用变量填充模板并返回/下载"""
    data = request.get_json() or {}
    template_id = data.get("id", "")
    variables = data.get("variables", {})

    template = None
    for t in BUILTIN_TEMPLATES:
        if t["id"] == template_id:
            template = t
            break

    if not template:
        return jsonify({"ok": False, "message": "模板不存在"}), 404

    content = template["content"]

    # 填充变量（未提供的保留原样）
    defaults = {
        "date": datetime.datetime.now().strftime("%Y-%m-%d"),
        "name": "", "dept": "", "project_name": "",
        "topic": "", "reason": "",
    }
    for k, v in defaults.items():
        if k not in variables:
            variables[k] = v
    for k, v in variables.items():
        content = content.replace("{" + k + "}", str(v))

    # 输出为 txt
    safe_name = re.sub(r'[\\/:*?"<>|]', '_', template["title"])
    filename = f"{safe_name}_{int(time.time())}.txt"
    filepath = OUTPUT_DIR / filename
    filepath.write_text(content, encoding="utf-8")

    return jsonify({
        "ok": True,
        "content": content,
        "filename": filename,
        "download_url": f"/api/download/{filename}",
    })


# ==================== FILE CONVERTER ====================
# 支持的转换类型
CONVERSION_CAPABILITIES = {
    "txt_to_docx": {
        "name": "TXT → DOCX (Word)",
        "input": "txt", "output": "docx",
        "desc": "纯文本转换为 Word 文档",
    },
    "docx_to_txt": {
        "name": "DOCX → TXT",
        "input": "docx", "output": "txt",
        "desc": "Word 文档提取纯文本",
    },
    "compress_image": {
        "name": "图片压缩",
        "input": "jpg/png/webp", "output": "jpg/png/webp",
        "desc": "压缩图片文件大小（可选择质量）",
    },
    "convert_image_format": {
        "name": "图片格式转换",
        "input": "jpg/png/webp/bmp/gif", "output": "多种",
        "desc": "转换图片格式（支持 jpg/png/webp/bmp 互转）",
    },

    # ===== 图片尺寸/裁剪/水印 =====
    "resize_image": {
        "name": "图片尺寸调整",
        "input": "jpg/png/webp/bmp", "output": "jpg/png/webp/bmp",
        "desc": "按指定宽高或百分比调整图片尺寸",
    },
    "crop_image": {
        "name": "图片裁剪",
        "input": "jpg/png/webp/bmp", "output": "jpg/png/webp/bmp",
        "desc": "按坐标区域裁剪图片",
    },
    "rotate_image": {
        "name": "图片旋转",
        "input": "jpg/png/webp/bmp", "output": "jpg/png/webp/bmp",
        "desc": "按角度旋转图片（90/180/270 或自定义角度）",
    },
    "watermark_image": {
        "name": "图片添加水印",
        "input": "jpg/png/webp/bmp", "output": "jpg/png/webp/bmp",
        "desc": "为图片添加文字水印",
    },
    "batch_convert_images": {
        "name": "批量图片处理",
        "input": "zip", "output": "zip",
        "desc": "批量调整图片尺寸/格式/质量（上传ZIP压缩包）",
    },
    "compress_image": {
        "name": "图片压缩",
        "input": "jpg/png/webp", "output": "jpg/png/webp",
        "desc": "压缩图片文件大小（可选择质量）",
    },
    "zip_compress": {
        "name": "文件压缩 (ZIP)",
        "input": "any", "output": "zip",
        "desc": "将文件或文件夹打包为 ZIP",
    },
    "zip_extract": {
        "name": "ZIP 解压",
        "input": "zip", "output": "folder",
        "desc": "解压 ZIP 压缩包",
    },
    "txt_encoding_convert": {
        "name": "TXT 编码转换",
        "input": "txt", "output": "txt",
        "desc": "转换文本文件编码（UTF-8/GBK/ASCII）",
    },

    # PDF/PPT/Word 互转
    "pdf_to_docx": {
        "name": "PDF → DOCX (Word)",
        "input": "pdf", "output": "docx",
        "desc": "PDF 文档转换为 Word 文档",
    },
    "docx_to_pdf": {
        "name": "DOCX → PDF",
        "input": "docx", "output": "pdf",
        "desc": "Word 文档转换为 PDF",
    },
    "pdf_to_pptx": {
        "name": "PDF → PPTX (PPT)",
        "input": "pdf", "output": "pptx",
        "desc": "PDF 文档逐页转为 PPT 幻灯片",
    },
    "pptx_to_pdf": {
        "name": "PPTX → PDF",
        "input": "pptx", "output": "pdf",
        "desc": "PPT 演示文稿转换为 PDF",
    },
    "docx_to_pptx": {
        "name": "DOCX → PPTX (PPT)",
        "input": "docx", "output": "pptx",
        "desc": "Word 文档大纲生成 PPT 演示文稿",
    },
    "pptx_to_docx": {
        "name": "PPTX → DOCX (Word)",
        "input": "pptx", "output": "docx",
        "desc": "PPT 内容提取为 Word 文档",
    },
}

@app.route("/api/convert/capabilities")
def api_convert_capabilities():
    return jsonify({"ok": True, "capabilities": CONVERSION_CAPABILITIES})

@app.route("/api/convert", methods=["POST"])
def api_convert():
    """执行文件转换"""
    if "file" not in request.files:
        return jsonify({"ok": False, "message": "没有上传文件"}), 400

    file = request.files["file"]
    convert_type = request.form.get("type", "")
    params_json = request.form.get("params", "{}")

    try:
        params = json.loads(params_json)
    except:
        params = {}

    if not convert_type or convert_type not in CONVERSION_CAPABILITIES:
        return jsonify({"ok": False, "message": "不支持的转换类型"}), 400

    cap = CONVERSION_CAPABILITIES[convert_type]
    orig_name = file.filename
    ext = os.path.splitext(orig_name)[1].lower().lstrip(".")

    # 检查输入格式
    if cap["input"] != "any" and ext not in cap["input"].split("/") and ext != cap["input"]:
        return jsonify({"ok": False, "message": f"文件格式 {ext} 不支持此转换"}), 400

    # 保存上传文件
    input_path = OUTPUT_DIR / f"upload_{int(time.time())}.{ext}"
    file.save(str(input_path))

    output_ext = cap["output"]
    if output_ext == "多种":
        output_ext = params.get("target_format", "jpg")
    if output_ext == "folder":
        output_ext = "zip"

    output_name = f"{os.path.splitext(orig_name)[0]}_converted.{output_ext}"
    output_path = OUTPUT_DIR / output_name

    try:
        if convert_type == "txt_to_docx":
            _convert_txt_to_docx(str(input_path), str(output_path))
        elif convert_type == "docx_to_txt":
            _convert_docx_to_txt(str(input_path), str(output_path))
        elif convert_type == "compress_image":
            quality = int(params.get("quality", 70))
            _compress_image(str(input_path), str(output_path), quality)
        elif convert_type == "convert_image_format":
            target = params.get("target_format", "jpg")
            _convert_image_format(str(input_path), str(output_path), target)
        elif convert_type == "zip_compress":
            _create_zip(str(input_path), str(output_path))
        elif convert_type == "zip_extract":
            out_dir = str(OUTPUT_DIR / f"extracted_{int(time.time())}")
            _extract_zip(str(input_path), out_dir)
            # zip output
            _create_zip(out_dir, str(output_path))
        elif convert_type == "txt_encoding_convert":
            target_enc = params.get("encoding", "utf-8")
            _convert_encoding(str(input_path), str(output_path), target_enc)
        elif convert_type == "resize_image":
            width = int(params.get("width", 0))
            height = int(params.get("height", 0))
            percent = float(params.get("percent", 0))
            _resize_image(str(input_path), str(output_path), width, height, percent)
        elif convert_type == "crop_image":
            left = int(params.get("left", 0))
            top = int(params.get("top", 0))
            right = int(params.get("right", 0))
            bottom = int(params.get("bottom", 0))
            _crop_image(str(input_path), str(output_path), left, top, right, bottom)
        elif convert_type == "rotate_image":
            angle = float(params.get("angle", 90))
            _rotate_image(str(input_path), str(output_path), angle)
        elif convert_type == "watermark_image":
            wm_text = params.get("text", "水印")
            position = params.get("position", "center")
            _watermark_image(str(input_path), str(output_path), wm_text, position)
        elif convert_type == "batch_convert_images":
            width_b = int(params.get("width", 0))
            height_b = int(params.get("height", 0))
            quality_b = int(params.get("quality", 85))
            target_fmt = params.get("target_format", "")
            _batch_process_images(str(input_path), str(output_path), width_b, height_b, quality_b, target_fmt)
        elif convert_type in ("pdf_to_docx", "docx_to_pdf", "pdf_to_pptx", "pptx_to_pdf", "docx_to_pptx", "pptx_to_docx"):
            _convert_office_doc(str(input_path), str(output_path), convert_type)
        else:
            return jsonify({"ok": False, "message": "未实现的转换"}), 400

        if not output_path.exists():
            return jsonify({"ok": False, "message": "转换失败"}), 500

        return jsonify({
            "ok": True,
            "output_file": output_name,
            "download_url": f"/api/download/{output_name}",
            "size": output_path.stat().st_size,
        })

    except Exception as e:
        return jsonify({"ok": False, "message": f"转换失败: {str(e)}"}), 500
    finally:
        # 清理上传文件
        if input_path.exists():
            input_path.unlink()


# ----- 转换实现 -----
def _convert_txt_to_docx(input_path, output_path):
    """TXT → DOCX 使用 python-docx"""
    try:
        from docx import Document
        from docx.shared import Pt, Inches
        doc = Document()
        doc.add_heading(os.path.basename(input_path), level=1)
        text = Path(input_path).read_text(encoding="utf-8")
        for line in text.split("\n"):
            if line.strip().startswith("===") or line.strip().startswith("=== "):
                doc.add_heading(line.strip("= "), level=2)
            elif line.strip():
                p = doc.add_paragraph(line)
        doc.save(output_path)
    except ImportError:
        # 降级：直接复制
        shutil.copy2(input_path, output_path)

def _convert_docx_to_txt(input_path, output_path):
    """DOCX → TXT"""
    try:
        from docx import Document
        doc = Document(input_path)
        text = "\n".join(p.text for p in doc.paragraphs)
        Path(output_path).write_text(text, encoding="utf-8")
    except ImportError:
        shutil.copy2(input_path, output_path)

def _compress_image(input_path, output_path, quality=70):
    """压缩图片"""
    try:
        from PIL import Image
        img = Image.open(input_path)
        fmt = img.format or "JPEG"
        save_kwargs = {"quality": quality, "optimize": True}
        if fmt == "PNG":
            save_kwargs = {"optimize": True}
        img.save(output_path, fmt, **save_kwargs)
    except ImportError:
        shutil.copy2(input_path, output_path)

def _convert_image_format(input_path, output_path, target_format):
    """图片格式转换"""
    try:
        from PIL import Image
        img = Image.open(input_path)
        # RGBA → RGB for JPEG
        if target_format.upper() == "JPEG" and img.mode == "RGBA":
            bg = Image.new("RGB", img.size, (255, 255, 255))
            bg.paste(img, mask=img.split()[3] if img.mode == "RGBA" else None)
            img = bg
        img.save(output_path, target_format.upper())
    except ImportError:
        shutil.copy2(input_path, output_path)

def _create_zip(input_path, output_path):
    """创建 ZIP 压缩"""
    with zipfile.ZipFile(output_path, "w", zipfile.ZIP_DEFLATED) as zf:
        if os.path.isdir(input_path):
            for root, dirs, files in os.walk(input_path):
                for f in files:
                    fp = os.path.join(root, f)
                    arcname = os.path.relpath(fp, os.path.dirname(input_path))
                    zf.write(fp, arcname)
        else:
            zf.write(input_path, os.path.basename(input_path))

def _extract_zip(input_path, output_dir):
    """解压 ZIP"""
    os.makedirs(output_dir, exist_ok=True)
    with zipfile.ZipFile(input_path, "r") as zf:
        zf.extractall(output_dir)

def _convert_encoding(input_path, output_path, target_encoding):
    """文本编码转换"""
    # 尝试检测原始编码
    raw = Path(input_path).read_bytes()
    for enc in ["utf-8", "gbk", "gb2312", "ascii", "utf-16"]:
        try:
            text = raw.decode(enc)
            break
        except:
            continue
    else:
        text = raw.decode("utf-8", errors="replace")

    Path(output_path).write_text(text, encoding=target_encoding, errors="replace")


def _convert_office_doc(input_path, output_path, convert_type):
    """PDF / PPT / Word 互转"""
    if convert_type == "pdf_to_docx":
        try:
            import fitz  # PyMuPDF
            from docx import Document
            from docx.shared import Inches, Pt
            doc = fitz.open(input_path)
            out = Document()
            for i, page in enumerate(doc):
                out.add_heading(f"第 {i+1} 页", level=2)
                text = page.get_text().strip()
                if text:
                    for line in text.split("\n"):
                        if line.strip():
                            out.add_paragraph(line.strip())
                # 提取图片
                for img_index, img_info in enumerate(page.get_images(full=True)):
                    try:
                        xref = img_info[0]
                        base_image = fitz.Pixmap(doc, xref)
                        if base_image.n - base_image.alpha > 3:
                            base_image = fitz.Pixmap(fitz.csRGB, base_image)
                        img_bytes = base_image.tobytes("png")
                        import io
                        out.add_picture(io.BytesIO(img_bytes), width=Inches(4.0))
                    except:
                        pass
            doc.close()
            out.save(output_path)
        except ImportError as e:
            raise ImportError(f"需要安装 PyMuPDF 和 python-docx: {e}")

    elif convert_type == "docx_to_pdf":
        try:
            from docx import Document
            import fitz
            doc = Document(input_path)
            # 用 fitz 创建空白 PDF，写入段落
            out_doc = fitz.open()
            page = out_doc.new_page()
            y_pos = 50
            for para in doc.paragraphs:
                if para.text.strip():
                    page.insert_text(fitz.Point(50, y_pos), para.text, fontsize=11)
                    y_pos += 20
                    if y_pos > 800:
                        page = out_doc.new_page()
                        y_pos = 50
            out_doc.save(output_path)
            out_doc.close()
        except ImportError as e:
            raise ImportError(f"需要安装 PyMuPDF 和 python-docx: {e}")

    elif convert_type == "pdf_to_pptx":
        try:
            import fitz
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            doc = fitz.open(input_path)
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            for i, page in enumerate(doc):
                slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
                text = page.get_text().strip()
                if text:
                    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(6.5))
                    tf = txBox.text_frame
                    tf.word_wrap = True
                    for j, line in enumerate(text.split("\n")):
                        if j == 0:
                            tf.text = line.strip()
                        elif line.strip():
                            p = tf.add_paragraph()
                            p.text = line.strip()
                # 加页号
                txBox = slide.shapes.add_textbox(Inches(11), Inches(7.0), Inches(2), Inches(0.4))
                tf = txBox.text_frame
                tf.text = f"{i+1}/{len(doc)}"
            doc.close()
            prs.save(output_path)
        except ImportError as e:
            raise ImportError(f"需要安装 PyMuPDF 和 python-pptx: {e}")

    elif convert_type == "pptx_to_pdf":
        try:
            from pptx import Presentation
            import fitz
            prs = Presentation(input_path)
            out_doc = fitz.open()
            slide_w = prs.slide_width or 9144000
            slide_h = prs.slide_height or 5143500
            for slide in prs.slides:
                rect = fitz.Rect(0, 0, slide_w * 96 / 914400, slide_h * 96 / 914400)
                page = out_doc.new_page(width=rect.width, height=rect.height)
                y = 50
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                page.insert_text(fitz.Point(50, y), text, fontsize=12)
                                y += 18
                                if y > rect.height - 50:
                                    break
            out_doc.save(output_path)
            out_doc.close()
        except ImportError as e:
            raise ImportError(f"需要安装 python-pptx 和 PyMuPDF: {e}")

    elif convert_type == "docx_to_pptx":
        try:
            from docx import Document
            from pptx import Presentation
            from pptx.util import Inches
            doc = Document(input_path)
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            current_slide = None
            for para in doc.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                style_name = para.style.name if para.style else ""
                if "Heading" in style_name or "标题" in style_name:
                    current_slide = prs.slides.add_slide(prs.slide_layouts[1])  # title+content
                    title_shape = current_slide.shapes.title
                    if title_shape:
                        title_shape.text = text
                elif current_slide:
                    body_shape = current_slide.placeholders[1] if len(current_slide.placeholders) > 1 else None
                    if body_shape:
                        tf = body_shape.text_frame
                        if not tf.text:
                            tf.text = text
                        else:
                            p = tf.add_paragraph()
                            p.text = text
                else:
                    current_slide = prs.slides.add_slide(prs.slide_layouts[6])
                    txBox = current_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
                    txBox.text_frame.text = text
            prs.save(output_path)
        except ImportError as e:
            raise ImportError(f"需要安装 python-docx 和 python-pptx: {e}")

    elif convert_type == "pptx_to_docx":
        try:
            from pptx import Presentation
            from docx import Document
            prs = Presentation(input_path)
            out = Document()
            out.add_heading("PPT 内容提取", level=1)
            for i, slide in enumerate(prs.slides):
                out.add_heading(f"幻灯片 {i+1}", level=2)
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                out.add_paragraph(text)
            out.save(output_path)
        except ImportError as e:
            raise ImportError(f"需要安装 python-pptx 和 python-docx: {e}")


# ==================== 图片处理函数 ====================
def _resize_image(input_path, output_path, width=0, height=0, percent=0):
    """调整图片尺寸"""
    try:
        from PIL import Image
        img = Image.open(input_path)
        if percent > 0:
            w = int(img.width * percent / 100)
            h = int(img.height * percent / 100)
        elif width > 0 and height > 0:
            w, h = width, height
        elif width > 0:
            ratio = width / img.width
            w, h = width, int(img.height * ratio)
        elif height > 0:
            ratio = height / img.height
            w, h = int(img.width * ratio), height
        else:
            return shutil.copy2(input_path, output_path)
        resized = img.resize((w, h), Image.LANCZOS)
        fmt = img.format or "JPEG"
        resized.save(output_path, fmt)
    except ImportError:
        shutil.copy2(input_path, output_path)


def _crop_image(input_path, output_path, left, top, right, bottom):
    """裁剪图片"""
    try:
        from PIL import Image
        img = Image.open(input_path)
        cropped = img.crop((left, top, right, bottom))
        fmt = img.format or "JPEG"
        cropped.save(output_path, fmt)
    except ImportError:
        shutil.copy2(input_path, output_path)


def _rotate_image(input_path, output_path, angle):
    """旋转图片"""
    try:
        from PIL import Image
        img = Image.open(input_path)
        rotated = img.rotate(angle, expand=True, resample=Image.BICUBIC)
        fmt = img.format or "JPEG"
        rotated.save(output_path, fmt)
    except ImportError:
        shutil.copy2(input_path, output_path)


def _watermark_image(input_path, output_path, text, position="center"):
    """添加文字水印"""
    try:
        from PIL import Image, ImageDraw, ImageFont
        img = Image.open(input_path).convert("RGBA")
        overlay = Image.new("RGBA", img.size, (255,255,255,0))
        draw = ImageDraw.Draw(overlay)
        font_size = max(20, min(img.width, img.height) // 15)
        try:
            font = ImageFont.truetype("simhei.ttf", font_size)
        except:
            font = ImageFont.load_default()
        bbox = draw.textbbox((0, 0), text, font=font)
        tw, th = bbox[2] - bbox[0], bbox[3] - bbox[1]
        margin = 20
        if position == "center":
            x, y = (img.width - tw) // 2, (img.height - th) // 2
        elif position == "topleft":
            x, y = margin, margin
        elif position == "topright":
            x, y = img.width - tw - margin, margin
        elif position == "bottomleft":
            x, y = margin, img.height - th - margin
        elif position == "bottomright":
            x, y = img.width - tw - margin, img.height - th - margin
        else:
            x, y = (img.width - tw) // 2, (img.height - th) // 2
        draw.text((x, y), text, font=font, fill=(255,255,255,160))
        result = Image.alpha_composite(img, overlay)
        result = result.convert("RGB")
        fmt = Image.open(input_path).format or "JPEG"
        result.save(output_path, fmt, quality=95)
    except ImportError:
        shutil.copy2(input_path, output_path)


def _batch_process_images(input_path, output_path, width=0, height=0, quality=85, target_fmt=""):
    """批量处理图片（ZIP包内所有图片）"""
    import tempfile, zipfile
    try:
        from PIL import Image
        temp_dir = Path(tempfile.mkdtemp())
        output_dir = temp_dir / "output"
        output_dir.mkdir(exist_ok=True)
        with zipfile.ZipFile(input_path, "r") as zf:
            zf.extractall(str(temp_dir / "input"))
        input_dir = temp_dir / "input"
        exts = {".jpg", ".jpeg", ".png", ".webp", ".bmp", ".gif", ".tiff"}
        count = 0
        for f in input_dir.rglob("*"):
            if f.suffix.lower() in exts:
                try:
                    img = Image.open(f)
                    if width > 0 and height > 0:
                        img = img.resize((width, height), Image.LANCZOS)
                    elif width > 0:
                        ratio = width / img.width
                        img = img.resize((width, int(img.height * ratio)), Image.LANCZOS)
                    elif height > 0:
                        ratio = height / img.height
                        img = img.resize((int(img.width * ratio), height), Image.LANCZOS)
                    out_name = f.stem
                    if target_fmt:
                        out_name += f".{target_fmt.lower()}"
                    else:
                        out_name += f.suffix
                    out_path = output_dir / out_name
                    save_kw = {"quality": quality, "optimize": True}
                    if target_fmt:
                        img.save(str(out_path), target_fmt.upper(), **save_kw)
                    else:
                        img.save(str(out_path), **save_kw)
                    count += 1
                except:
                    pass
        # 打包结果
        with zipfile.ZipFile(output_path, "w", zipfile.ZIP_DEFLATED) as zf:
            for f in output_dir.rglob("*"):
                arcname = f.relative_to(output_dir)
                zf.write(f, arcname)
        # 清理
        import shutil as _sh
        _sh.rmtree(temp_dir, ignore_errors=True)
    except ImportError:
        shutil.copy2(input_path, output_path)


# ==================== DOWNLOAD & STATIC ====================
@app.route("/api/download/<filename>")
def api_download(filename):
    safe = os.path.basename(filename)
    return send_file(str(OUTPUT_DIR / safe), as_attachment=True,
                     download_name=safe)

@app.route("/")
def index():
    return app.send_static_file("index.html")


# ==================== STARTUP ====================



if __name__ == "__main__":
    # 打印依赖检查
    missing = []
    try:
        from PIL import Image
        print("[OK] Pillow → 图片处理")
    except ImportError:
        missing.append("pillow")
        print("[WARN] Pillow 未安装 → 图片转换降级为复制")

    try:
        from docx import Document
        print("[OK] python-docx → Word 文档处理")
    except ImportError:
        missing.append("python-docx")
        print("[WARN] python-docx 未安装 → TXT/DOCX 转换降级为复制")

    try:
        import fitz
        print("[OK] PyMuPDF → PDF 处理")
    except ImportError:
        missing.append("PyMuPDF")
        print("[WARN] PyMuPDF 未安装 → PDF 转换降级")

    try:
        from pptx import Presentation
        print("[OK] python-pptx → PPT 处理")
    except ImportError:
        missing.append("python-pptx")
        print("[WARN] python-pptx 未安装 → PPT 转换降级")

    print("=" * 50)
    print("  Office Toolkit v1.0")
    print("  办公文件工具箱")
    print("=" * 50)
    print(f"  [模板] {len(BUILTIN_TEMPLATES)} 个内置模板")
    print(f"  [转换] {len(CONVERSION_CAPABILITIES)} 种转换类型")
    if missing:
        print(f"\n  缺失依赖 (功能受限):")
        for m in missing:
            print(f"     pip install --trusted-host pypi.org --trusted-host files.pythonhosted.org {m}")
    print(f"\n  打开: http://127.0.0.1:18700")
    print("  Ctrl+C 停止\n")

    import webbrowser
    import threading as _th
    _th.Timer(0.8, lambda: webbrowser.open("http://127.0.0.1:18700")).start()

    try:
        app.run(host="127.0.0.1", port=18700, debug=False, threaded=True)
    finally:
        pass
